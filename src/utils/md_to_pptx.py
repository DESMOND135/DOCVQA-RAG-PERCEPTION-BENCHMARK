import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from PIL import Image

import matplotlib.pyplot as plt
from matplotlib import mathtext

# Professional Academic Palette
BLUE_NAVY = RGBColor(0, 51, 102) # #003366
BLUE_ARXIV = RGBColor(31, 119, 180) # #1f77b4
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

def render_latex_to_image(latex, out_path):
    """Renders a LaTeX string into a high-resolution transparent PNG in academic STIX serif font."""
    import matplotlib.pyplot as plt
    plt.rcParams['mathtext.fontset'] = 'stix'
    plt.rcParams['font.family'] = 'STIXGeneral'
    
    fig, ax = plt.subplots(figsize=(6, 1.2))
    fig.patch.set_facecolor('none')
    ax.axis('off')
    
    # Center the math equations perfectly
    ax.text(
        0.5, 0.5, f"${latex}$",
        ha='center', va='center', fontsize=26, color='#000000',
        transform=ax.transAxes
    )
    
    plt.savefig(out_path, transparent=True, bbox_inches='tight', pad_inches=0.02, dpi=300)
    plt.close(fig)

def get_image_dimensions(img_path):
    try:
        with Image.open(img_path) as im:
            return im.size
    except Exception as e:
        print(f"PIL size reading error: {e}")
        return (800, 600)

def style_title_slide(slide, title_text, subtitle_text):
    # Absolute theme consistency: White background with Dark Blue Headings
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = WHITE
    
    # Perfect horizontal and vertical centralization on widescreen 16:9 canvas
    title_shape = slide.shapes.title
    title_shape.text = title_text
    title_shape.left = Inches(0.5)
    title_shape.width = Inches(12.333)
    title_shape.top = Inches(1.6)
    title_shape.height = Inches(2.2)
    
    for p in title_shape.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = BLUE_NAVY
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.name = 'Calibri'
    
    # Strip any markdown asterisks from the subtitle text for the title slide
    clean_subtitle = subtitle_text.replace('**', '').replace('*', '')
    subtitle = slide.placeholders[1]
    subtitle.text = clean_subtitle
    subtitle.left = Inches(0.5)
    subtitle.width = Inches(12.333)
    subtitle.top = Inches(3.8)
    subtitle.height = Inches(3.2)
    
    for p in subtitle.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = BLACK
        p.font.size = Pt(16)
        p.font.name = 'Calibri'

def add_slide_decorations(slide, current, total):
    """Adds branding and slide numbers with professional Navy accent."""
    if current > 0:
        # Progress bar at bottom
        bar_width = (current / total) * 13.333
        bar = slide.shapes.add_shape(6, 0, Inches(7.4), Inches(bar_width), Inches(0.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = BLUE_NAVY; bar.line.fill.background()
        
        box = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.8), Inches(0.4))
        p = box.text_frame.paragraphs[0]; p.text = f"{current} / {total}"
        p.font.size = Pt(10); p.font.color.rgb = BLUE_NAVY; p.alignment = PP_ALIGN.RIGHT

def generate_defense_deck(md_path, pptx_path):
    print(f"Generating Presentation Deck: {md_path} -> {pptx_path}")
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    with open(md_path, 'r', encoding='utf-8') as f: content = f.read()
    slides_raw = [s for s in content.split('\n---\n') if s.strip()]

    for i, slide_md in enumerate(slides_raw):
        lines = [l.strip() for l in slide_md.split('\n') if l.strip()]
        if not lines: continue
        
        title, bullets, img_path, table_data, math_images = "Document Section", [], None, [], []
        img_caption, table_caption = None, None
        if lines[0].startswith('#'):
            title = lines[0].lstrip('#').strip(); body_lines = lines[1:]
        else:
            body_lines = lines

        if i == 0 or "Thank You" in title:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            style_title_slide(slide, title, "\n".join(body_lines))
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        add_slide_decorations(slide, i, len(slides_raw)-1)
        
        # Remove default content placeholder if present to allow custom textbox grid placement
        for shape in list(slide.shapes):
            if shape.is_placeholder and shape.placeholder_format.idx == 1:
                sp = shape._element
                sp.getparent().remove(sp)
        
        # Header - ALWAYS CENTRALIZED
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.left = Inches(0.5)
        title_shape.width = Inches(12.333)
        title_shape.top = Inches(0.4)
        title_shape.height = Inches(0.8)
        
        title_p = title_shape.text_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(28)
        title_p.font.color.rgb = BLUE_NAVY
        title_p.font.bold = True
        title_p.font.name = 'Calibri'
        
        for line in body_lines:
            if line.startswith('- ') or line.startswith('* '):
                if '$$' not in line:
                    bullets.append(line[2:].strip())
            elif line.startswith('|'):
                if '---' not in line:
                    cols = [c.strip() for c in line.split('|') if c.strip()]
                    if cols: table_data.append(cols)
            elif '$$' in line:
                m = re.search(r'\$\$(.*?)\$\$', line)
                if m:
                    latex = m.group(1).strip()
                    os.makedirs('temp_math', exist_ok=True)
                    math_img = os.path.join('temp_math', f'math_{i}_{len(math_images)}.png')
                    render_latex_to_image(latex, math_img)
                    math_images.append(math_img)
            elif '![' in line:
                m = re.search(r'!\[.*?\]\((.*?)\)', line); 
                if m:
                    rel_path = m.group(1)
                    p = os.path.normpath(os.path.join(os.path.dirname(md_path), rel_path))
                    
                    # Robust resolution to handle differences between 1-level and 2-level parent setups
                    if not os.path.exists(p):
                        if rel_path.startswith('../'):
                            p = os.path.normpath(os.path.join(os.path.dirname(md_path), '../' + rel_path))
                    
                    # Absolute fallback to the workspace root folders
                    if not os.path.exists(p):
                        workspace_root = 'c:/Users/Administrator/Downloads/THESIS PROJECT'
                        fixed_rel = rel_path.replace('../', '').replace('..\\', '')
                        p = os.path.normpath(os.path.join(workspace_root, fixed_rel))
                        
                    if os.path.exists(p):
                        img_path = p
                        print(f"Presentation loaded image: {img_path}")
                    else:
                        print(f"Presentation image NOT found: {rel_path}")
            elif line.startswith('**Figure') or (line.startswith('**') and 'Figure' in line):
                img_caption = line.replace('**', '').strip()
            elif line.startswith('**Table') or (line.startswith('**') and 'Table' in line):
                table_caption = line.replace('**', '').strip()

        # Layout Dispatching based on components
        if table_data:
            # Table Slide Layout
            if bullets:
                txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.8))
                tf = txBox.text_frame
                tf.word_wrap = True
                for b_index, b_text in enumerate(bullets):
                    p = tf.add_paragraph() if b_index > 0 else tf.paragraphs[0]
                    p.space_after = Pt(6)
                    p.line_spacing = 1.2
                    run = p.add_run()
                    run.text = b_text.replace('**', '').replace('*', '')
                    run.font.size = Pt(15)
                    run.font.name = 'Calibri'
                    run.font.italic = True
                    run.font.color.rgb = BLUE_ARXIV
                    p.alignment = PP_ALIGN.CENTER
            
            rows, cols = len(table_data), len(table_data[0])
            t_width = Inches(11.733)
            t_left = Inches(0.8)
            t_top = Inches(2.1) if bullets else Inches(1.5)
            t_height = Inches(4.5) if bullets else Inches(5.2)
            
            if table_caption:
                t_top = Inches(2.0) if bullets else Inches(1.4)
                t_height = Inches(4.4) if bullets else Inches(5.0)
            
            table = slide.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height).table
            
            if table_caption:
                caption_top = t_top + t_height + Inches(0.1)
                captionBox = slide.shapes.add_textbox(t_left, caption_top, t_width, Inches(0.4))
                cp = captionBox.text_frame.paragraphs[0]
                cp.text = table_caption
                cp.font.size = Pt(13)
                cp.font.name = 'Calibri'
                cp.font.bold = True
                cp.font.color.rgb = BLUE_NAVY
                cp.alignment = PP_ALIGN.CENTER
            for r in range(rows):
                for c in range(cols):
                    cell = table.cell(r, c)
                    val = table_data[r][c].replace('**', '').replace('*', '')
                    cell.text = val
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(13) if cols > 7 else Pt(15)
                    p.font.name = 'Calibri'
                    p.font.bold = (r == 0 or "Hybrid" in val)
                    p.alignment = PP_ALIGN.CENTER
                    
                    if r == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = BLUE_NAVY
                        p.font.color.rgb = WHITE
                    else:
                        cell.fill.solid()
                        if r % 2 == 0:
                            cell.fill.fore_color.rgb = RGBColor(242, 246, 250)
                        else:
                            cell.fill.fore_color.rgb = WHITE
                            
                        if "Hybrid" in val or "Proposed" in val:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(220, 235, 248)
                            p.font.color.rgb = BLUE_NAVY

        elif math_images and bullets:
            # Bullet Explanations (Left) + Multi-Math Formulas (Right) Split Layout
            bullets_left = Inches(0.5)
            bullets_width = Inches(5.8)
            bullets_top = Inches(1.4)
            bullets_height = Inches(5.4)
            
            txBox = slide.shapes.add_textbox(bullets_left, bullets_top, bullets_width, bullets_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            for b_index, b_text in enumerate(bullets):
                p = tf.add_paragraph() if b_index > 0 else tf.paragraphs[0]
                p.space_after = Pt(20) # Vertical spacing between sections
                p.level = 0
                p.line_spacing = 1.2 # Line spacing for premium readability
                
                parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', b_text)
                for part in parts:
                    run = p.add_run()
                    if part.startswith('**'): run.text = part[2:-2]; run.font.bold = True
                    elif part.startswith('*'): run.text = part[1:-1]; run.font.italic = True
                    else: run.text = part
                    run.font.size = Pt(19) # Larger body font size
                    run.font.name = 'Calibri'
                    run.font.color.rgb = BLACK
            
            img_max_w = Inches(6.2)
            img_max_h = Inches(5.4)
            
            # Widen visual columns and balance spatial occupancy
            rem_left = bullets_left + bullets_width + Inches(0.3)
            rem_width = prs.slide_width - rem_left - Inches(0.5)
            
            if len(math_images) == 1:
                math_path = math_images[0]
                orig_w, orig_h = get_image_dimensions(math_path)
                aspect = orig_w / orig_h
                w = img_max_w
                h = w / aspect
                if h > Inches(3.6):
                    h = Inches(3.6)
                    w = h * aspect
                top = Inches(1.4) + (img_max_h - h) / 2
                left = rem_left + (rem_width - w) / 2
                slide.shapes.add_picture(math_path, left, top, width=w, height=h)
            elif len(math_images) == 2:
                for idx, math_path in enumerate(math_images):
                    orig_w, orig_h = get_image_dimensions(math_path)
                    aspect = orig_w / orig_h
                    w = img_max_w
                    h = w / aspect
                    if h > Inches(2.2):
                        h = Inches(2.2)
                        w = h * aspect
                    left = rem_left + (rem_width - w) / 2
                    top = Inches(1.4) + idx * Inches(2.6)
                    slide.shapes.add_picture(math_path, left, top, width=w, height=h)
            else:
                # 4 equations stacked in a beautiful high-density vertical layout
                for idx, math_path in enumerate(math_images):
                    orig_w, orig_h = get_image_dimensions(math_path)
                    aspect = orig_w / orig_h
                    w = img_max_w
                    h = w / aspect
                    if h > Inches(1.0):
                        h = Inches(1.0)
                        w = h * aspect
                    left = rem_left + (rem_width - w) / 2
                    top = Inches(1.4) + idx * Inches(1.3)
                    slide.shapes.add_picture(math_path, left, top, width=w, height=h)

        elif img_path and bullets:
            if "Core of Intelligence" in title or "Data at the Core" in title:
                # Top-and-Bottom Stacked Layout for Widescreen Widespan Infographics
                text_left = Inches(0.8)
                text_width = Inches(11.733)
                text_top = Inches(1.3)
                text_height = Inches(1.2)
                
                txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                tf = txBox.text_frame
                tf.word_wrap = True
                
                for b_index, b_text in enumerate(bullets):
                    p = tf.add_paragraph() if b_index > 0 else tf.paragraphs[0]
                    p.space_after = Pt(4)
                    p.level = 0
                    p.line_spacing = 1.1
                    p.alignment = PP_ALIGN.CENTER
                    
                    parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', b_text)
                    for part in parts:
                        run = p.add_run()
                        if part.startswith('**'): run.text = part[2:-2]; run.font.bold = True
                        elif part.startswith('*'): run.text = part[1:-1]; run.font.italic = True
                        else: run.text = part
                        run.font.size = Pt(15) # Centered top callout font size
                        run.font.name = 'Calibri'
                        run.font.color.rgb = BLACK
                
                # Widen and maximize landscape diagram size
                img_max_w = Inches(11.733)
                img_max_h = Inches(4.3)
                if img_caption:
                    img_max_h = Inches(3.9)
                    
                orig_w, orig_h = get_image_dimensions(img_path)
                aspect = orig_w / orig_h
                
                w = img_max_w
                h = w / aspect
                if h > img_max_h:
                    h = img_max_h
                    w = h * aspect
                    
                img_left = (prs.slide_width - w) / 2
                img_top = Inches(2.5) + (img_max_h - h) / 2
                
                slide.shapes.add_picture(img_path, img_left, img_top, width=w, height=h)
                
                if img_caption:
                    captionBox = slide.shapes.add_textbox(img_left, img_top + h + Inches(0.02), w, Inches(0.3))
                    cp = captionBox.text_frame.paragraphs[0]
                    cp.text = img_caption
                    cp.font.size = Pt(11)
                    cp.font.name = 'Calibri'
                    cp.font.bold = True
                    cp.font.color.rgb = BLUE_NAVY
                    cp.alignment = PP_ALIGN.CENTER
            else:
                # Bullet Explanations (Left) + Diagram (Right) Split Layout
                bullets_left = Inches(0.5)
                bullets_width = Inches(5.8)
                bullets_top = Inches(1.4)
                bullets_height = Inches(5.4)
                
                txBox = slide.shapes.add_textbox(bullets_left, bullets_top, bullets_width, bullets_height)
                tf = txBox.text_frame
                tf.word_wrap = True
                
                for b_index, b_text in enumerate(bullets):
                    p = tf.add_paragraph() if b_index > 0 else tf.paragraphs[0]
                    p.space_after = Pt(18) # Widen vertical spacing between bullets
                    p.level = 0
                    p.line_spacing = 1.2 # Line spacing for premium readability
                    
                    parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', b_text)
                    for part in parts:
                        run = p.add_run()
                        if part.startswith('**'): run.text = part[2:-2]; run.font.bold = True
                        elif part.startswith('*'): run.text = part[1:-1]; run.font.italic = True
                        else: run.text = part
                        run.font.size = Pt(19) # Larger body font size
                        run.font.name = 'Calibri'
                        run.font.color.rgb = BLACK
                
                # Widen visual columns and balance spatial occupancy
                rem_left = bullets_left + bullets_width + Inches(0.3)
                rem_width = prs.slide_width - rem_left - Inches(0.5)
                
                # Substantially enlarged maximum frame to utilize margins and balance better
                img_max_w = Inches(6.2)
                img_max_h = Inches(5.4)
                if img_caption:
                    img_max_h = Inches(4.7)
                    
                orig_w, orig_h = get_image_dimensions(img_path)
                aspect = orig_w / orig_h
                
                w = img_max_w
                h = w / aspect
                if h > img_max_h:
                    h = img_max_h
                    w = h * aspect
                    
                img_left = rem_left + (rem_width - w) / 2
                img_top = Inches(1.4) # Aligned at the very top of content area to match text block
                
                slide.shapes.add_picture(img_path, img_left, img_top, width=w, height=h)
                
                if img_caption:
                    captionBox = slide.shapes.add_textbox(img_left, img_top + h + Inches(0.1), w, Inches(0.4))
                    cp = captionBox.text_frame.paragraphs[0]
                    cp.text = img_caption
                    cp.font.size = Pt(12)
                    cp.font.name = 'Calibri'
                    cp.font.bold = True
                    cp.font.color.rgb = BLUE_NAVY
                    cp.alignment = PP_ALIGN.CENTER

        elif img_path or math_images:
            # Centered Layout for Single Image/Equation (No Bullets)
            single_path = img_path if img_path else math_images[0]
            img_max_w = Inches(12.5) # Enlarged to fill slide dimensions perfectly
            img_max_h = Inches(5.9) # Substantially increased to let landscape images scale up
            if img_caption:
                img_max_h = Inches(5.5) # Optimized to leave exact space for caption and title
                
            orig_w, orig_h = get_image_dimensions(single_path)
            aspect = orig_w / orig_h
            
            w = img_max_w
            h = w / aspect
            if h > img_max_h:
                h = img_max_h
                w = h * aspect
                
            img_left = (prs.slide_width - w) / 2
            img_top = Inches(1.2) + (img_max_h - h) / 2
            
            slide.shapes.add_picture(single_path, img_left, img_top, width=w, height=h)
            
            if img_caption:
                captionBox = slide.shapes.add_textbox(img_left, img_top + h + Inches(0.05), w, Inches(0.35))
                cp = captionBox.text_frame.paragraphs[0]
                cp.text = img_caption
                cp.font.size = Pt(11)
                cp.font.name = 'Calibri'
                cp.font.bold = True
                cp.font.color.rgb = BLUE_NAVY
                cp.alignment = PP_ALIGN.CENTER

        else:
            # Standard Text Slide Layout (Bullets only - Text-heavy Slides)
            bullets_left = Inches(0.8)
            bullets_width = Inches(11.733)
            bullets_top = Inches(1.4)
            bullets_height = Inches(5.4)
            
            txBox = slide.shapes.add_textbox(bullets_left, bullets_top, bullets_width, bullets_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            for b_index, b_text in enumerate(bullets):
                p = tf.add_paragraph() if b_index > 0 else tf.paragraphs[0]
                p.space_after = Pt(24) # Increased paragraph spacing to occupy vertical space beautifully
                p.level = 0
                p.line_spacing = 1.2 # Premium academic line spacing
                
                parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', b_text)
                for part in parts:
                    run = p.add_run()
                    if part.startswith('**'): run.text = part[2:-2]; run.font.bold = True
                    elif part.startswith('*'): run.text = part[1:-1]; run.font.italic = True
                    else: run.text = part
                    run.font.size = Pt(22) # Moderately larger body text for absolute slide fullness
                    run.font.name = 'Calibri'
                    run.font.color.rgb = BLACK

    prs.save(pptx_path)
    print(f"Presentation Generated: {pptx_path}")

if __name__ == "__main__":
    # Robust path resolution
    md_path = 'MAIN/presentation/presentation.md'
    pptx_path = 'MAIN/presentation/presentation.pptx'
    if not os.path.exists(md_path):
        # Fallback if run from THESIS directory directly
        md_path = 'presentation/presentation.md'
        pptx_path = 'presentation/presentation.pptx'
    generate_defense_deck(md_path, pptx_path)
